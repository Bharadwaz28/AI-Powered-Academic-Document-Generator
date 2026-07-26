from pptx import Presentation
from pptx.util import Inches, Pt
import os
import requests
from io import BytesIO
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def generate_ppt(slides, output_file="output.pptx"):

    # Create outputs folder
    output_dir = "outputs"
    os.makedirs(output_dir, exist_ok=True)
    file_path = os.path.join(output_dir, output_file)

    prs = Presentation()

    # ---------------- MAIN SLIDES ----------------
    for idx, slide_data in enumerate(slides):

        layout_type = idx % 3

        if layout_type == 0:

            slide = prs.slides.add_slide(prs.slide_layouts[1])

            slide.shapes.title.text = slide_data.get("title", "No Title")

            content = slide.placeholders[1]

            content.left = Inches(0.5)
            content.top = Inches(1.9)
            content.width = Inches(5.5)
            content.height = Inches(5)

            tf = content.text_frame
            tf.clear()
            tf.word_wrap = True

            bullets = slide_data.get("bullets", [])[:5]

            for i, bullet in enumerate(bullets):

                if i == 0:
                    p = tf.paragraphs[0]

                else:
                    p = tf.add_paragraph()

                p.text = bullet
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(18)

                p.space_after = Pt(18)
                p.line_spacing = 1.2

            image_url = slide_data.get("image_url")

            if isinstance(image_url, str) and image_url.startswith("http"):

                try:

                    response = requests.get(image_url, timeout=5)

                    if response.status_code == 200:

                        slide.shapes.add_picture(
                            BytesIO(response.content),
                            Inches(6.0),
                            Inches(1.9),
                            width=Inches(3.5),
                            height=Inches(5)
                        )

                except Exception as e:
                    print(e)

        elif layout_type == 1:

            slide = prs.slides.add_slide(prs.slide_layouts[1])

            slide.shapes.title.text = slide_data.get("title", "No Title")

            image_url = slide_data.get("image_url")

            if isinstance(image_url, str) and image_url.startswith("http"):

                try:

                    response = requests.get(image_url, timeout=5)

                    if response.status_code == 200:

                        slide.shapes.add_picture(
                            BytesIO(response.content),
                            Inches(1.6),
                            Inches(1.66),
                            width=Inches(6.8),
                            height=Inches(2.9)
                        )

                except Exception as e:
                    print(e)

            content = slide.placeholders[1]

            content.left = Inches(0.5)
            content.top = Inches(4.7)
            content.width = Inches(9)
            content.height = Inches(2.2)

            tf = content.text_frame
            tf.clear()

            tf.word_wrap = True

            bullets = slide_data.get("bullets", [])[:4]

            for i, bullet in enumerate(bullets):

                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = bullet
                p.level = 0

                for run in p.runs:
                    run.font.size = Pt(18)

                p.space_after = Pt(18)
                p.line_spacing = 1

        else:

            slide = prs.slides.add_slide(prs.slide_layouts[1])

            slide.shapes.title.text = slide_data.get("title", "No Title")

            content = slide.placeholders[1]

            content.left = Inches(0.5)
            content.top = Inches(1.66)
            content.width = Inches(9)
            content.height = Inches(2.2)

            tf = content.text_frame
            tf.clear()

            tf.word_wrap = True

            bullets = slide_data.get("bullets", [])[:4]

            for i, bullet in enumerate(bullets):

                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = bullet
                p.level = 0

                for run in p.runs:
                    run.font.size = Pt(18)

                p.space_after = Pt(18)
                p.line_spacing = 1

            image_url = slide_data.get("image_url")

            if isinstance(image_url, str) and image_url.startswith("http"):

                try:

                    response = requests.get(image_url, timeout=5)

                    if response.status_code == 200:

                        slide.shapes.add_picture(
                            BytesIO(response.content),
                            Inches(1.6),
                            Inches(4),
                            width=Inches(6.8),
                            height=Inches(2.9)
                        )

                except Exception as e:
                    print(e)

        # ================= SLIDE 2: TABLE =================
        if slide_data.get("table_data"):
            table_data = slide_data["table_data"]

            table_slide = prs.slides.add_slide(prs.slide_layouts[5])

            # ---- TABLE TITLE ----
            if table_slide.shapes.title:
                table_slide.shapes.title.text = table_data.get(
                    "table_title",
                    slide_data.get("title", "No Title") + " Table"
                )

            rows = len(table_data["rows"]) + 1
            cols = len(table_data["headers"])

            table = table_slide.shapes.add_table(
                rows, cols,
                Inches(0.5), Inches(1.9),
                Inches(9), Inches(5)
            ).table

            # Headers
            for col_idx, header in enumerate(table_data["headers"]):
                table.cell(0, col_idx).text = header

            # Data
            for row_idx, row in enumerate(table_data["rows"], start=1):
                for col_idx, value in enumerate(row):
                    table.cell(row_idx, col_idx).text = str(value)

        # ================= SLIDE 3: CHART =================
        if slide_data.get("chart_data"):

            chart = slide_data["chart_data"]

            labels = chart.get("labels", [])
            values = chart.get("values", [])

            valid = (
                chart.get("chart_title")
                and chart.get("metric")
                and len(labels) == len(values)
                and len(labels) >= 2
            )

            if not valid:
                continue

            chart_slide = prs.slides.add_slide(prs.slide_layouts[5])

            # ---- CHART TITLE ----
            if chart_slide.shapes.title:
                chart_slide.shapes.title.text = chart.get(
                    "chart_title",
                    slide_data.get("title", "No Title") + " Chart"
                )

            chart_type = chart.get("type", "").lower()

            labels = chart.get("labels", [])
            values = chart.get("values", [])

            # ================= PIE CHART  =================
            if chart_type == "pie":

                clean_labels = []
                clean_values = []

                for l, v in zip(labels, values):
                    try:
                        v = float(v)
                        if v > 0:
                            clean_labels.append(l)
                            clean_values.append(v)
                    except:
                        continue

                if len(clean_values) < 2:
                    print("Skipping pie chart: not enough valid data")
                    continue

                labels = clean_labels
                values = clean_values

            # ================= OTHER CHARTS =================
            else:
                clean_labels = []
                clean_values = []

                for l, v in zip(labels, values):
                    try:
                        v = float(v)
                        clean_labels.append(l)
                        clean_values.append(v)
                    except:
                        continue

                if len(clean_values) < 1:
                    print("Skipping chart: no valid data")
                    continue

                labels = clean_labels
                values = clean_values

            chart_data = CategoryChartData()
            chart_data.categories = labels

            if chart_type == "pie":
                chart_data.add_series("", values)
            else:
                chart_data.add_series(
                    chart.get("metric", "Data"),
                    values
                )
            if chart_type == "line":
                chart_enum = XL_CHART_TYPE.LINE
            elif chart_type == "pie":
                chart_enum = XL_CHART_TYPE.PIE
            else:
                chart_enum = XL_CHART_TYPE.COLUMN_CLUSTERED

            # ---- ADD CHART ----
            chart_shape = chart_slide.shapes.add_chart(
                chart_enum,
                Inches(0.5), Inches(1.9),
                Inches(9), Inches(5),
                chart_data
            )

            chart_obj = chart_shape.chart

            plot = chart_obj.plots[0]
            plot.has_data_labels = True

            data_labels = plot.data_labels
            data_labels.show_value = False

            # PIE CHART LABELS
            if chart_type == "pie":
                data_labels.show_category_name = True
                data_labels.show_percentage = True
            else:
                data_labels.show_category_name = True
                data_labels.show_value = True

            # ================= AXIS TITLES =================
            if chart_type != "pie":
                chart_obj.category_axis.has_title = True
                chart_obj.category_axis.axis_title.text_frame.text = chart.get(
                    "x_label", "X-Axis")

                chart_obj.value_axis.has_title = True
                chart_obj.value_axis.axis_title.text_frame.text = chart.get(
                    "y_label",
                    chart.get("metric", "Value")
                )

    # Save file
    prs.save(file_path)

    return file_path
